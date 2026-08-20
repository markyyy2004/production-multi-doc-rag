import os
import tempfile
import pytest
import numpy as np
import docx
from pptx import Presentation
from app import extract_docs_from_file, safe_ocr_extract

class DummyProgressBar:
    def progress(self, val): pass

class DummyStatusText:
    def text(self, val): pass

@pytest.fixture
def dummy_ui():
    return DummyProgressBar(), DummyStatusText()

def test_extract_from_empty_and_corrupted_files(dummy_ui):
    pbar, status = dummy_ui
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
        tmp_path = tmp.name
    try:
        assert extract_docs_from_file(tmp_path, "empty.txt", pbar, status) == []
    finally:
        if os.path.exists(tmp_path): os.remove(tmp_path)

def test_extract_docx_clean_tokens(dummy_ui):
    pbar, status = dummy_ui
    doc = docx.Document()
    doc.add_paragraph("Production RAG Pipeline Architecture")
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp_path = tmp.name
    doc.save(tmp_path)
    try:
        docs = extract_docs_from_file(tmp_path, "sample.docx", pbar, status)
        assert len(docs) == 1
        assert "Production RAG Pipeline Architecture" in docs[0].page_content
    finally:
        if os.path.exists(tmp_path): os.remove(tmp_path)

def test_extract_pptx_clean_tokens(dummy_ui):
    pbar, status = dummy_ui
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Microservices Concurrency"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
    prs.save(tmp_path)
    try:
        docs = extract_docs_from_file(tmp_path, "sample.pptx", pbar, status)
        assert len(docs) == 1
        assert "Microservices Concurrency" in docs[0].page_content
        assert docs[0].metadata["slide"] == 1
    finally:
        if os.path.exists(tmp_path): os.remove(tmp_path)

def test_ocr_blank_scan_fallback():
    blank_white = np.ones((150, 150, 3), dtype=np.uint8) * 255
    text = safe_ocr_extract(blank_white)
    assert text == ""