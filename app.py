import io
import os
import shutil
import pickle
import logging
import concurrent.futures
import docx
import fitz  # PyMuPDF
import numpy as np
from PIL import Image
from pptx import Presentation
from rapidocr_onnxruntime import RapidOCR
import streamlit as st
from dotenv import load_dotenv

# LangChain Core LCEL & Groq
from langchain_community.vectorstores import FAISS
from langchain_community.retrievers import BM25Retriever
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage, AIMessage
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables import RunnableLambda
from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Database Layer
from database import (
    init_db,
    register_user,
    verify_user,
    save_chat_message,
    load_user_chat_history,
    clear_user_chat_history
)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

load_dotenv()
init_db()

DB_INDEX_PATH = "faiss_index"
BM25_INDEX_PATH = "bm25_retriever.pkl"
TEMP_DOCS_DIR = "uploaded_docs_cache"
os.makedirs(TEMP_DOCS_DIR, exist_ok=True)

MAX_USER_QUERY_CHARS = 1500
OCR_CONFIDENCE_THRESHOLD = 0.50

# --- PAGE CONFIG ---
st.set_page_config(
    page_title="Nexus RAG // Intelligent Knowledge Engine",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- DARK GLASSMORPHISM UI ---
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800&family=JetBrains+Mono:wght@400;500&display=swap');

    .stApp {
        background-color: #080A0F;
        color: #E2E8F0;
        font-family: 'Plus Jakarta Sans', sans-serif;
    }

    /* Ambient background glow */
    .stApp::before {
        content: '';
        position: fixed;
        top: -15%;
        left: 50%;
        transform: translateX(-50%);
        width: 900px;
        height: 450px;
        background: radial-gradient(circle, rgba(56, 189, 248, 0.08) 0%, rgba(139, 92, 246, 0.05) 50%, rgba(0,0,0,0) 80%);
        z-index: 0;
        pointer-events: none;
    }

    /* Modern Glassmorphic Sidebar */
    [data-testid="stSidebar"] {
        background-color: #0D1117 !important;
        border-right: 1px solid rgba(255, 255, 255, 0.08);
        box-shadow: 10px 0 35px rgba(0,0,0,0.5);
    }

    /* Chat message bubble styling */
    [data-testid="stChatMessage"] {
        background: rgba(15, 23, 42, 0.4) !important;
        border: 1px solid rgba(255, 255, 255, 0.05) !important;
        border-radius: 12px !important;
        margin-bottom: 12px !important;
        padding: 14px 18px !important;
        backdrop-filter: blur(8px);
    }

    /* Quick suggestion buttons */
    .stButton > button {
        background: #111827 !important;
        color: #F8FAFC !important;
        border: 1px solid rgba(255, 255, 255, 0.12) !important;
        border-radius: 8px !important;
        font-size: 0.82rem !important;
        font-weight: 600 !important;
        transition: all 0.2s cubic-bezier(0.4, 0, 0.2, 1) !important;
    }

    .stButton > button:hover {
        background: #1E293B !important;
        border-color: #38BDF8 !important;
        color: #38BDF8 !important;
        transform: translateY(-2px);
        box-shadow: 0 4px 14px rgba(56, 189, 248, 0.2) !important;
    }

    /* Inline code blocks */
    code {
        background: rgba(56, 189, 248, 0.12) !important;
        color: #38BDF8 !important;
        padding: 2px 6px !important;
        border-radius: 4px !important;
        font-family: 'JetBrains Mono', monospace !important;
    }
</style>
""", unsafe_allow_html=True)


# --- 1. MODEL INITIALIZATION ---
@st.cache_resource
def load_llm():
    return ChatGroq(
        model="openai/gpt-oss-20b",
        temperature=0,
        groq_api_key=os.getenv("GROQ_API_KEY"),
    )


@st.cache_resource
def load_embeddings():
    return HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        encode_kwargs={"batch_size": 64}
    )


@st.cache_resource
def load_ocr():
    return RapidOCR()


llm = load_llm()
embeddings = load_embeddings()
ocr_engine = load_ocr()


# --- 2. HARDENED HYBRID RAG PIPELINE (RRF) ---
def create_hybrid_retriever(vector_store, all_chunks=None):
    faiss_retriever = vector_store.as_retriever(search_kwargs={"k": 10})

    bm25_retriever = None
    if all_chunks:
        try:
            bm25_retriever = BM25Retriever.from_documents(all_chunks)
            bm25_retriever.k = 10
            with open(BM25_INDEX_PATH, "wb") as f:
                pickle.dump(bm25_retriever, f)
        except Exception as e:
            logger.error(f"Failed to build BM25 retriever: {e}")
            bm25_retriever = None
    elif os.path.exists(BM25_INDEX_PATH):
        try:
            with open(BM25_INDEX_PATH, "rb") as f:
                bm25_retriever = pickle.load(f)
                bm25_retriever.k = 10
        except Exception as e:
            logger.error(f"Failed to load BM25 retriever: {e}")
            bm25_retriever = None

    def hybrid_search(query: str, k: int = 2, w_dense: float = 0.6, w_sparse: float = 0.4):
        dense_docs = []
        sparse_docs = []

        try:
            dense_docs = faiss_retriever.invoke(query)
        except Exception as e:
            logger.error(f"FAISS search failed: {e}")

        if bm25_retriever:
            query_terms = set(query.lower().split())
            try:
                raw_sparse = bm25_retriever.invoke(query)
                sparse_docs = [
                    d for d in raw_sparse 
                    if any(term in d.page_content.lower() for term in query_terms)
                ]
            except Exception as e:
                logger.error(f"BM25 search failed: {e}")

        doc_scores = {}
        doc_map = {}

        for rank, doc in enumerate(dense_docs):
            content = doc.page_content
            doc_map[content] = doc
            doc_scores[content] = doc_scores.get(content, 0.0) + (w_dense / (60 + rank + 1))

        for rank, doc in enumerate(sparse_docs):
            content = doc.page_content
            doc_map[content] = doc
            doc_scores[content] = doc_scores.get(content, 0.0) + (w_sparse / (60 + rank + 1))

        sorted_contents = sorted(doc_scores.keys(), key=lambda c: doc_scores[c], reverse=True)
        return [doc_map[c] for c in sorted_contents[:k]]

    return RunnableLambda(hybrid_search)


def build_conversational_rag_chain(hybrid_retriever):
    rephrase_system_prompt = (
        "Rephrase the follow-up question into a concise standalone query based on chat history. "
        "Do NOT answer, return only the standalone question."
    )
    rephrase_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", rephrase_system_prompt),
            MessagesPlaceholder(variable_name="chat_history"),
            ("human", "{question}"),
        ]
    )
    question_rephraser = rephrase_prompt | llm | StrOutputParser()

    qa_system_prompt = (
        "You are Nexus, an AI technical expert. Use the retrieved context to answer the question. "
        "Be direct, accurate, and concise. Support LaTeX math ($...$). If not covered, state it clearly.\n\n"
        "Context:\n{context}"
    )
    qa_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", qa_system_prompt),
            MessagesPlaceholder(variable_name="chat_history"),
            ("human", "{question}"),
        ]
    )
    qa_chain = qa_prompt | llm | StrOutputParser()

    def get_sources_and_stream(inputs):
        raw_query = inputs.get("question", "").strip()
        safe_query = raw_query[:MAX_USER_QUERY_CHARS]
        trimmed_history = inputs.get("chat_history", [])[-2:] if inputs.get("chat_history") else []

        if trimmed_history:
            try:
                standalone_query = question_rephraser.invoke({
                    "question": safe_query,
                    "chat_history": trimmed_history
                })
            except Exception as e:
                logger.warning(f"Rephraser fallback: {e}")
                standalone_query = safe_query
        else:
            standalone_query = safe_query

        try:
            retrieved_docs = hybrid_retriever.invoke(standalone_query)
        except Exception as e:
            logger.error(f"Retrieval error: {e}")
            retrieved_docs = []

        if not retrieved_docs:
            def empty_stream():
                yield "I searched the indexed documents, but could not find any relevant information matching your query."
            return empty_stream(), []

        context_str = "\n\n".join(doc.page_content for doc in retrieved_docs)

        try:
            stream_gen = qa_chain.stream({
                "context": context_str,
                "chat_history": trimmed_history,
                "question": safe_query,
            })
            return stream_gen, retrieved_docs
        except Exception as e:
            logger.error(f"LLM generation failed: {e}")
            def error_stream():
                yield f"An error occurred during response generation: {str(e)}"
            return error_stream(), retrieved_docs

    return get_sources_and_stream


# --- 3. HARDENED DOCUMENT & OCR EXTRACTOR ---
def safe_ocr_extract(img_np):
    try:
        result, _ = ocr_engine(img_np)
        if not result:
            return ""
        
        valid_lines = []
        for item in result:
            if len(item) >= 3:
                text, score = item[1], item[2]
                if score >= OCR_CONFIDENCE_THRESHOLD and text.strip():
                    valid_lines.append(text.strip())
            elif len(item) >= 2 and item[1].strip():
                valid_lines.append(item[1].strip())

        return "\n".join(valid_lines)
    except Exception as e:
        logger.warning(f"OCR execution warning: {e}")
        return ""


def extract_docs_from_file(file_path, filename, progress_bar, status_text):
    ext = os.path.splitext(filename)[1].lower()
    page_docs = []

    if not os.path.exists(file_path) or os.path.getsize(file_path) == 0:
        logger.warning(f"Skipping empty or missing file: {filename}")
        return []

    # Parallel Multi-Threaded PDF Processing
    if ext == ".pdf":
        doc = None
        try:
            doc = fitz.open(file_path)
            total_pages = len(doc)

            def process_single_page(page_idx):
                try:
                    page = doc.load_page(page_idx)
                    page_text = page.get_text()

                    if page_text and len(page_text.strip()) > 40:
                        return (page_idx, Document(
                            page_content=page_text.strip(),
                            metadata={"source": filename, "page": page_idx + 1, "filepath": file_path}
                        ))
                    else:
                        pix = page.get_pixmap(dpi=150)
                        img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
                        ocr_text = safe_ocr_extract(np.array(img))

                        if ocr_text:
                            return (page_idx, Document(
                                page_content=ocr_text,
                                metadata={"source": filename, "page": page_idx + 1, "type": "OCR", "filepath": file_path}
                            ))
                except Exception as e:
                    logger.warning(f"Error processing page {page_idx + 1} of {filename}: {e}")
                return (page_idx, None)

            status_text.text(f"Indexing '{filename}' ({total_pages} pages) in parallel...")
            extracted_results = []
            completed_count = 0

            with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
                futures = [executor.submit(process_single_page, i) for i in range(total_pages)]
                for future in concurrent.futures.as_completed(futures):
                    result = future.result()
                    if result and result[1] is not None:
                        extracted_results.append(result)
                    completed_count += 1
                    progress_bar.progress(completed_count / max(total_pages, 1))

            extracted_results.sort(key=lambda x: x[0])
            page_docs = [doc_item[1] for doc_item in extracted_results]
        except Exception as e:
            logger.error(f"Corrupted PDF file {filename}: {e}")
            st.warning(f"⚠️ Could not parse '{filename}'. The file may be corrupt or encrypted.")
        finally:
            if doc:
                doc.close()

    elif ext == ".docx":
        try:
            status_text.text(f"Indexing Word Document: {filename}...")
            doc = docx.Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs if p.text and p.text.strip()])
            if text.strip():
                page_docs.append(Document(page_content=text.strip(), metadata={"source": filename, "filepath": file_path}))
        except Exception as e:
            logger.error(f"Corrupted DOCX {filename}: {e}")
            st.warning(f"⚠️ Could not parse Word document '{filename}'.")

    elif ext == ".pptx":
        try:
            status_text.text(f"Indexing Slides: {filename}...")
            prs = Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                cleaned_text = "\n".join(slide_text).strip()
                if cleaned_text:
                    page_docs.append(Document(
                        page_content=cleaned_text,
                        metadata={"source": filename, "slide": idx + 1, "filepath": file_path}
                    ))
        except Exception as e:
            logger.error(f"Corrupted PPTX {filename}: {e}")
            st.warning(f"⚠️ Could not parse PowerPoint '{filename}'.")

    elif ext in [".png", ".jpg", ".jpeg"]:
        try:
            status_text.text(f"Running RapidOCR on image: {filename}...")
            ocr_text = safe_ocr_extract(file_path)
            if ocr_text:
                page_docs.append(Document(
                    page_content=ocr_text,
                    metadata={"source": filename, "type": "OCR", "filepath": file_path}
                ))
        except Exception as e:
            logger.error(f"Image OCR failed on {filename}: {e}")

    elif ext in [".txt", ".csv", ".py", ".java", ".json", ".md", ".cpp", ".c"]:
        try:
            status_text.text(f"Reading source code: {filename}...")
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read().strip()
                if text:
                    page_docs.append(Document(page_content=text, metadata={"source": filename, "filepath": file_path}))
        except Exception as e:
            logger.error(f"Failed to read file {filename}: {e}")

    return page_docs


def cache_uploaded_files(uploaded_files):
    file_info = []
    for uploaded_file in uploaded_files:
        try:
            save_path = os.path.join(TEMP_DOCS_DIR, uploaded_file.name)
            with open(save_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            file_info.append((save_path, uploaded_file.name))
        except Exception as e:
            logger.error(f"Failed to save temp file {uploaded_file.name}: {e}")
    return file_info


def load_and_split_documents(file_info_list, progress_bar, status_text):
    raw_docs = []
    for save_path, original_name in file_info_list:
        docs = extract_docs_from_file(save_path, original_name, progress_bar, status_text)
        raw_docs.extend(docs)

    if not raw_docs:
        return []

    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=80)
    return splitter.split_documents(raw_docs)


# --- 4. AUTHENTICATION VIEW ---
def auth_view():
    st.markdown("""
    <div style="text-align: center; padding: 40px 0 20px 0;">
        <span style="background: rgba(56, 189, 248, 0.1); border: 1px solid rgba(56, 189, 248, 0.25); color: #38BDF8; padding: 4px 14px; border-radius: 9999px; font-size: 0.75rem; font-weight: 700; text-transform: uppercase;">
            🔐 Multi-User Security
        </span>
        <h2 style="font-size: 2.2rem; margin-top: 10px;">Nexus Knowledge Copilot</h2>
        <p style="color: #94A3B8;">Sign in to access your private indexed workspace.</p>
    </div>
    """, unsafe_allow_html=True)

    col1, col2, col3 = st.columns([1, 1.2, 1])
    with col2:
        tab_choice = st.radio("Mode", ["Login", "Register"], horizontal=True, label_visibility="collapsed")
        
        with st.form("auth_form"):
            username = st.text_input("Username").strip()
            password = st.text_input("Password", type="password")
            submit = st.form_submit_button("Proceed", use_container_width=True)

            if submit:
                if not username or not password:
                    st.error("Please enter both username and password.")
                    return

                if tab_choice == "Register":
                    if register_user(username, password):
                        st.success("Account created successfully! Please switch to Login.")
                    else:
                        st.error("Username is already taken.")
                else:
                    user = verify_user(username, password)
                    if user:
                        st.session_state.authenticated = True
                        st.session_state.user = user
                        
                        db_messages = load_user_chat_history(user["id"])
                        st.session_state.messages = []
                        st.session_state.chat_history = []
                        for m in db_messages:
                            st.session_state.messages.append(m)
                            if m["role"] == "user":
                                st.session_state.chat_history.append(HumanMessage(content=m["content"]))
                            else:
                                st.session_state.chat_history.append(AIMessage(content=m["content"]))
                        st.rerun()
                    else:
                        st.error("Invalid username or password.")


# --- 5. MAIN APPLICATION CONTROLLER ---
def main():
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False

    if not st.session_state.authenticated:
        auth_view()
        st.stop()

    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "rag_chain" not in st.session_state:
        st.session_state.rag_chain = None
    if "preview_target" not in st.session_state:
        st.session_state.preview_target = None
    if "input_prompt_val" not in st.session_state:
        st.session_state.input_prompt_val = ""

    if st.session_state.rag_chain is None and os.path.exists(DB_INDEX_PATH):
        try:
            saved_vectorstore = FAISS.load_local(
                DB_INDEX_PATH, 
                embeddings, 
                allow_dangerous_deserialization=True
            )
            hybrid_ret = create_hybrid_retriever(saved_vectorstore)
            st.session_state.rag_chain = build_conversational_rag_chain(hybrid_ret)
        except Exception as e:
            logger.error(f"Could not load persisted FAISS index: {e}")
            st.session_state.rag_chain = None

    # --- SIDEBAR CONTROLS ---
    with st.sidebar:
        st.markdown(f"👤 Logged in as: **`{st.session_state.user['username']}`**")
        if st.button("🚪 Log Out", use_container_width=True):
            st.session_state.authenticated = False
            st.session_state.user = None
            st.session_state.messages = []
            st.session_state.chat_history = []
            st.session_state.preview_target = None
            st.rerun()

        st.markdown("---")
        st.markdown("### ⚡ NEXUS ENGINE")
        st.caption("Hybrid BM25 + FAISS Vector Engine")

        st.markdown("#### 💾 Hybrid Index Status")
        if os.path.exists(DB_INDEX_PATH) and st.session_state.rag_chain is not None:
            st.success("🟢 Hybrid Index Active (BM25 + FAISS)")
        else:
            st.info("🟡 Awaiting Document Indexing")

        st.markdown("---")
        st.markdown("#### 📂 Document Hub")
        uploaded_files = st.file_uploader(
            "Upload files",
            type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg", "py", "java", "json", "csv"],
            accept_multiple_files=True,
            label_visibility="collapsed"
        )

        if st.button("⚡ Process & Index Documents", use_container_width=True) and uploaded_files:
            progress_bar = st.progress(0)
            status_text = st.empty()

            file_info = cache_uploaded_files(uploaded_files)
            chunks = load_and_split_documents(file_info, progress_bar, status_text)

            if chunks:
                status_text.text("Building FAISS & BM25 indexes...")
                try:
                    vector_store = FAISS.from_documents(chunks, embeddings)
                    vector_store.save_local(DB_INDEX_PATH)
                    
                    hybrid_ret = create_hybrid_retriever(vector_store, all_chunks=chunks)
                    st.session_state.rag_chain = build_conversational_rag_chain(hybrid_ret)

                    progress_bar.empty()
                    status_text.empty()
                    st.success(f"Indexed {len(chunks)} chunks into Hybrid Retrieval Engine!")
                    st.rerun()
                except Exception as e:
                    progress_bar.empty()
                    status_text.empty()
                    st.error(f"Indexing error: {e}")
            else:
                progress_bar.empty()
                status_text.empty()
                st.error("No valid text could be extracted from uploaded files.")

        st.markdown("---")
        st.markdown("#### ⚙️ Operations")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("🧹 Clear Chat", use_container_width=True):
                clear_user_chat_history(st.session_state.user["id"])
                st.session_state.messages = []
                st.session_state.chat_history = []
                st.session_state.preview_target = None
                st.rerun()
        with col2:
            if st.button("🗑️ Wipe DB", use_container_width=True):
                if os.path.exists(DB_INDEX_PATH):
                    shutil.rmtree(DB_INDEX_PATH)
                if os.path.exists(BM25_INDEX_PATH):
                    os.remove(BM25_INDEX_PATH)
                if os.path.exists(TEMP_DOCS_DIR):
                    shutil.rmtree(TEMP_DOCS_DIR)
                    os.makedirs(TEMP_DOCS_DIR, exist_ok=True)
                st.session_state.rag_chain = None
                clear_user_chat_history(st.session_state.user["id"])
                st.session_state.messages = []
                st.session_state.chat_history = []
                st.session_state.preview_target = None
                st.rerun()

    # --- MAIN VIEW: SPLIT-SCREEN LAYOUT ---
    if st.session_state.preview_target:
        chat_col, preview_col = st.columns([1.15, 0.85], gap="medium")
    else:
        chat_col = st.container()
        preview_col = None

    with chat_col:
        if len(st.session_state.messages) == 0:
            st.markdown("""
            <div style="text-align: center; padding: 30px 10px 10px 10px;">
                <span style="background: rgba(56, 189, 248, 0.1); border: 1px solid rgba(56, 189, 248, 0.25); color: #38BDF8; padding: 4px 14px; border-radius: 9999px; font-size: 0.75rem; font-weight: 700; text-transform: uppercase;">
                    ⚡ Hybrid Semantic & Exact Retrieval
                </span>
                <h1 style="font-size: 2.5rem; margin-top: 12px; margin-bottom: 6px;">Nexus Knowledge Copilot</h1>
                <p style="color: #94A3B8; font-size: 0.95rem; max-width: 650px; margin: 0 auto 20px auto;">
                    Exact technical keyword precision (BM25) combined with dense conceptual understanding (FAISS).
                </p>
            </div>
            """, unsafe_allow_html=True)

        for msg_idx, msg in enumerate(st.session_state.messages):
            with st.chat_message(msg["role"], avatar="⚡" if msg["role"] == "assistant" else "👤"):
                st.markdown(msg["content"])
                
                # Check for sources in memory or loaded from DB
                sources = msg.get("sources", [])
                if not sources and "sources_json" in msg:
                    sources = [Document(page_content=s["page_content"], metadata=s["metadata"]) for s in msg["sources_json"]]

                if msg["role"] == "assistant" and sources:
                    with st.expander("🔍 Verified Document Citations"):
                        for doc_idx, doc in enumerate(sources):
                            meta = doc.metadata
                            src = meta.get("source", "Unknown")
                            page_num = meta.get("page", None)
                            slide_num = meta.get("slide", None)
                            
                            c1, c2 = st.columns([0.75, 0.25])
                            with c1:
                                loc = f" (Page {page_num})" if page_num else (f" (Slide {slide_num})" if slide_num else "")
                                st.markdown(f"**Source {doc_idx + 1}:** `{src}`{loc}")
                                st.caption(doc.page_content[:200] + ("..." if len(doc.page_content) > 200 else ""))
                            with c2:
                                if st.button("👁️ Inspect", key=f"btn_src_{msg_idx}_{doc_idx}"):
                                    st.session_state.preview_target = {
                                        "source": src,
                                        "page": page_num,
                                        "content": doc.page_content,
                                        "filepath": meta.get("filepath", os.path.join(TEMP_DOCS_DIR, src))
                                    }
                                    st.rerun()

        # Quick Suggestions
        st.markdown("<div style='height: 10px;'></div>", unsafe_allow_html=True)
        q1, q2, q3 = st.columns(3)
        with q1:
            if st.button("💡 Summarize Core Concepts", use_container_width=True):
                st.session_state.input_prompt_val = "Summarize core concepts and fundamental architecture covered across the notes."
        with q2:
            if st.button("⚔️ Compare Key Differences", use_container_width=True):
                st.session_state.input_prompt_val = "Compare key differences and syntax patterns across the notes."
        with q3:
            if st.button("📝 Generate Practice Quiz", use_container_width=True):
                st.session_state.input_prompt_val = "Generate a 3-question technical quiz based on the notes."

        # Submission Handler
        query_val = st.chat_input("Ask anything about your documents, code, or diagrams...") or st.session_state.input_prompt_val
        if st.session_state.input_prompt_val:
            st.session_state.input_prompt_val = ""

        if query_val:
            if not st.session_state.rag_chain:
                st.warning("⚠️ Please upload and index documents via the sidebar first!")
                return

            save_chat_message(st.session_state.user["id"], "user", query_val)
            st.session_state.messages.append({"role": "user", "content": query_val})
            with st.chat_message("user", avatar="👤"):
                st.markdown(query_val)

            with st.chat_message("assistant", avatar="⚡"):
                stream_gen, sources = st.session_state.rag_chain({
                    "question": query_val,
                    "chat_history": st.session_state.chat_history,
                })
                
                full_response = st.write_stream(stream_gen)

                if sources:
                    with st.expander("🔍 Verified Document Citations"):
                        for doc_idx, doc in enumerate(sources):
                            meta = doc.metadata
                            src = meta.get("source", "Unknown")
                            page_num = meta.get("page", None)
                            slide_num = meta.get("slide", None)
                            
                            c1, c2 = st.columns([0.75, 0.25])
                            with c1:
                                loc = f" (Page {page_num})" if page_num else (f" (Slide {slide_num})" if slide_num else "")
                                st.markdown(f"**Source {doc_idx + 1}:** `{src}`{loc}")
                                st.caption(doc.page_content[:200] + ("..." if len(doc.page_content) > 200 else ""))
                            with c2:
                                if st.button("👁️ Inspect", key=f"btn_src_new_{doc_idx}"):
                                    st.session_state.preview_target = {
                                        "source": src,
                                        "page": page_num,
                                        "content": doc.page_content,
                                        "filepath": meta.get("filepath", os.path.join(TEMP_DOCS_DIR, src))
                                    }
                                    st.rerun()

            save_chat_message(st.session_state.user["id"], "assistant", full_response, sources=sources)
            st.session_state.chat_history.append(HumanMessage(content=query_val))
            st.session_state.chat_history.append(AIMessage(content=full_response))
            st.session_state.messages.append({
                "role": "assistant",
                "content": full_response,
                "sources": sources
            })
            st.rerun()

    # --- SPLIT-SCREEN INSPECTOR PANEL ---
    if preview_col and st.session_state.preview_target:
        with preview_col:
            st.markdown("""
            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 12px;">
                <h4 style="margin: 0; color: #38BDF8;">📄 Document Inspector</h4>
            </div>
            """, unsafe_allow_html=True)

            target = st.session_state.preview_target
            file_path = target.get("filepath", "")
            page_num = target.get("page", 1)

            st.markdown(f"**File:** `{target['source']}`")
            if page_num:
                st.caption(f"Viewing Page: {page_num}")

            if st.button("❌ Close Inspector", use_container_width=True):
                st.session_state.preview_target = None
                st.rerun()

            if os.path.exists(file_path) and file_path.lower().endswith(".pdf") and page_num:
                doc = None
                try:
                    doc = fitz.open(file_path)
                    page = doc.load_page(page_num - 1)
                    pix = page.get_pixmap(dpi=150)
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    st.image(img, caption=f"{target['source']} - Page {page_num}", use_container_width=True)
                except Exception as e:
                    logger.warning(f"Could not render image: {e}")
                    st.text_area("Extracted Text", target["content"], height=400)
                finally:
                    if doc:
                        doc.close()
            else:
                st.text_area("Indexed Text Chunk", target["content"], height=450)


if __name__ == "__main__":
    main()